"""Build the final project slide deck (.pptx) from verified results.

Every number here is copied from the README result tables / results CSVs — nothing is
computed or invented at build time. Figures are embedded from assets/.

Run:  python scripts/build_presentation.py
Output: docs/Robustness_of_Vision_Models_to_Distortions.pptx
"""
from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "assets"
OUT = ROOT / "docs" / "Robustness_of_Vision_Models_to_Distortions.pptx"

NAVY = RGBColor(0x1F, 0x38, 0x64)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x70, 0x70, 0x70)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

REPO = "github.com/yonatan-gif/Image-Processing-Final-Project"


def _set(tf, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    r.font.name = "Calibri"


def title_bar(slide, title):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.45); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    _set(tf, title, 25, WHITE, bold=True)
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.0), W, Inches(0.06))
    accent.fill.solid(); accent.fill.fore_color.rgb = BLUE; accent.line.fill.background()
    accent.shadow.inherit = False


def footer(slide, n):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.35))
    tf = box.text_frame
    _set(tf, f"{REPO}", 9, GREY)
    pn = slide.shapes.add_textbox(Inches(12.4), Inches(7.05), Inches(0.8), Inches(0.35))
    _set(pn.text_frame, str(n), 9, GREY, align=PP_ALIGN.RIGHT)


def bullets(slide, items, left, top, width, height, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        txt, lvl = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ("   " * lvl) + ("•  " if lvl == 0 else "–  ") + txt
        p.space_after = Pt(8)
        p.line_spacing = 1.05
        for r in p.runs:
            r.font.size = Pt(size - 2 * lvl)
            r.font.color.rgb = DARK
            r.font.name = "Calibri"
    return box


def fit(image_path, max_w_in, max_h_in):
    w, h = Image.open(image_path).size
    scale = min(max_w_in / w, max_h_in / h)
    return Inches(w * scale), Inches(h * scale)


def figure(slide, image_path, top_in, max_w_in=12.3, max_h_in=None, bottom_in=7.0, region_left_in=None):
    if max_h_in is None:
        max_h_in = bottom_in - top_in
    iw, ih = fit(str(image_path), max_w_in, max_h_in)
    if region_left_in is None:  # center across the whole slide
        left = (13.333 - iw.inches) / 2
    else:  # center within the right region [region_left_in, 13.1]
        left = region_left_in + max(0.0, (13.1 - region_left_in - iw.inches) / 2)
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top_in), width=iw, height=ih)


_n = [0]


def new_slide(title):
    slide = prs.slides.add_slide(BLANK)
    title_bar(slide, title)
    _n[0] += 1
    footer(slide, _n[0])
    return slide


# ---------------------------------------------------------------- Slide 1: title
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.55), W, Inches(0.06))
strip.fill.solid(); strip.fill.fore_color.rgb = BLUE; strip.line.fill.background(); strip.shadow.inherit = False
t = s.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.0))
tf = t.text_frame; tf.word_wrap = True
_set(tf, "Robustness of Vision Models to Image Distortions", 40, WHITE, bold=True)
sub = s.shapes.add_textbox(Inches(0.95), Inches(4.8), Inches(11.5), Inches(2.0))
tf = sub.text_frame; tf.word_wrap = True
for i, line in enumerate([
    "Image-Processing & Computer-Vision — Final Project",
    "Oxford-IIIT Pet  ·  SIFT  ·  ResNet-50  ·  DeepLabV3",
    REPO,
]):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    r = p.runs[0]; r.font.size = Pt(18 if i < 2 else 15); r.font.color.rgb = WHITE; r.font.name = "Calibri"
    p.space_after = Pt(6)
_n[0] = 1

# ---------------------------------------------------------------- Slide 2: objective
s = new_slide("Objective & Central Question")
bullets(s, [
    "Measure how three common distortions — Gaussian noise, blur, JPEG compression — degrade vision algorithms.",
    "Two levels of abstraction: low-level SIFT keypoints vs. high-level classification & segmentation.",
    "For every condition, compare two recovery strategies: clean the image (restoration) vs. adapt the model (fine-tuning).",
    "Central question: does cleaning the image or adapting the model recover more task performance?",
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=20)

# ---------------------------------------------------------------- Slide 3: bundle
s = new_slide("The End-to-End Bundle")
bullets(s, [
    "Dataset: Oxford-IIIT Pet — 3,680 images, 37 breeds; ground truth for classification and segmentation.",
    "Task 1 (high-level, DL): Classification — ResNet-50 — Top-1 accuracy.",
    "Task 2 (high-level, dense DL): Segmentation — DeepLabV3-ResNet50 — mIoU.",
    "Task 3 (low-level, classical): Interest points — SIFT — repeatability + matching score.",
    "Distortions: Gaussian noise (σ 5–80), Gaussian blur (σ 0.5–8), JPEG (quality 90–10).",
    "Recovery: classical restoration (NLM / unsharp / bilateral) and fine-tuning on distorted data (DL tasks).",
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=18)

# ---------------------------------------------------------------- Slide 4: dataset/EDA
s = new_slide("Dataset & EDA — Oxford-IIIT Pet")
bullets(s, [
    "3,680 images · 37 breeds (~100 each, balanced) · median size 500×375.",
    "Ground truth: breed label, head bounding box, trimap segmentation mask.",
    "Below: breed label (title) + head bbox (green) + pet mask (red).",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.6), size=16)
figure(s, ASSETS / "eda_samples.png", top_in=3.1, bottom_in=6.9)

# ---------------------------------------------------------------- Slide 5: distortions + SNR
s = new_slide("Distortions and Measured SNR")
bullets(s, [
    "Each intensity level mapped to a physical SNR (mean PSNR over 40 images).",
    "Noise: 34 → 12 dB   ·   Blur: 39 → 22 dB   ·   JPEG: 39 → 27 dB.",
    "Rows below: clean → weak → mid → strong → restored (matched cleaner).",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.6), size=16)
figure(s, ASSETS / "distortion_grid.png", top_in=3.1, bottom_in=6.9)

# ---------------------------------------------------------------- Slide 6: experiment design
s = new_slide("Experiment Design & Metrics")
bullets(s, [
    "Four conditions vs. the clean baseline:  distorted → restored → fine-tuned.",
    "Recovery Δ = metric(recovered) − metric(distorted)   (positive = helps, negative = hurts).",
    "Metrics: Top-1 (classification); mIoU + per-class IoU (segmentation); repeatability + matching (SIFT).",
    "Eval sets: 300 val images (DL aggregate) · 2,480 held-out images (per-breed) · 40 seeded images, mean ± std (SIFT).",
    "Fixed seed (42) throughout for reproducibility.",
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=18)

# ---------------------------------------------------------------- Slide 7: Task 1
s = new_slide("Task 1 — Classification (ResNet-50)")
bullets(s, [
    "Clean baseline Top-1 = 0.933 (300 val images).",
    "Top-1  distorted / restored / fine-tuned:",
    ("noise σ=5:  0.93 / 0.70 / 0.90        noise σ=80:  0.22 / 0.24 / 0.75", 1),
    ("blur σ=4:  0.44 / 0.50 / 0.79         JPEG q=10:  0.64 / 0.64 / 0.82", 1),
    "Robust to mild distortion; blind restoration mostly hurts; fine-tuning is the best recovery.",
], Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.4), size=16)
figure(s, ASSETS / "classification_finetune_gaussian_noise.png", top_in=1.7, max_w_in=6.2, max_h_in=4.9, region_left_in=6.8)

# ---------------------------------------------------------------- Slide 8: per-breed
s = new_slide("Task 1 — Per-Breed × Intensity")
bullets(s, [
    "Per-breed from 2,480 held-out images (~67/breed); clean overall Top-1 = 0.911.",
    "Weakest breeds are lookalike pairs: Egyptian Mau 0.72, Am. Pit Bull 0.74, Staffordshire 0.74, Birman 0.79; three breeds at 1.00.",
    "Moderate distortion amplifies existing confusions; extreme noise (σ=80) flattens all breeds.",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.7), size=16)
figure(s, ASSETS / "classification_per_class.png", top_in=3.3, bottom_in=6.9)

# ---------------------------------------------------------------- Slide 9: Task 2
s = new_slide("Task 2 — Segmentation (DeepLabV3)")
bullets(s, [
    "Clean baseline mIoU = 0.923 — the most robust of the three tasks.",
    "mIoU  distorted / restored / fine-tuned:",
    ("noise σ=80:  0.62 / 0.62 / 0.87     blur σ=8:  0.79 / 0.80 / 0.83     JPEG q=10:  0.85 / 0.87 / 0.89", 1),
    "Per-class IoU (noise σ=80): background 0.73 vs. pet 0.52 (gap 0.22); fine-tuning closes it to 0.04.",
    "Degradation concentrates in the textured pet, not the smooth background.",
], Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.4), size=15)
figure(s, ASSETS / "segmentation_gaussian_noise.png", top_in=1.7, max_w_in=6.2, max_h_in=4.9, region_left_in=6.8)

# ---------------------------------------------------------------- Slide 10: Task 3
s = new_slide("Task 3 — SIFT Keypoints")
bullets(s, [
    "40 seeded images (mean ± std); ~714 keypoints per clean image.",
    "Repeatability  distorted → restored (Δ):",
    ("noise σ=5:  0.77 → 0.42  (−0.35)      blur σ=1:  0.45 → 0.76  (+0.31)      JPEG q=90:  0.86 → 0.54  (−0.32)", 1),
    "Blur erases keypoints (290 → 12); noise / JPEG spawn spurious ones (→ 362 / 331).",
    "The cleaner helps only blur; smoothing removes the fine texture keypoints sit on.",
], Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.9), size=16)
figure(s, ASSETS / "keypoints_visual.png", top_in=3.5, bottom_in=6.9)

# ---------------------------------------------------------------- Slide 11: fairness setup
s = new_slide("Fairness Check — Was “Restoration Hurts” a Dose Artifact?")
bullets(s, [
    "All cleaners ran at ONE fixed strength across a 16× noise range (NLM h=10): overdosed at σ=5, underdosed at σ=80.",
    "Blind protocol: estimate noise σ̂ from the image (Immerkær 1996, validated — monotone, small spread), set h = 0.8·σ̂.",
    "Same pixels-only interface as the fine-tuned model — a fair, deployment-realistic opponent.",
], Inches(0.6), Inches(1.4), Inches(6.0), Inches(5.4), size=15)
figure(s, ASSETS / "classification_noise_fairness.png", top_in=1.7, max_w_in=6.2, max_h_in=4.9, region_left_in=6.8)

# ---------------------------------------------------------------- Slide 12: fairness results
s = new_slide("Fairness Check — Results")
bullets(s, [
    "Overdosing was real: classification σ=5 damage −0.22 (fixed) → −0.07 (blind); SIFT −0.35 → −0.10.",
    "Sign holds for texture tasks: correctly-dosed denoising still loses to the noisy input (classification up to σ=40; SIFT at every level).",
    "Segmentation flips: blind restoration at σ=80 recovers mIoU 0.62 → 0.82 (fine-tune 0.87) — a genuine dose artifact there.",
    "Cross-task law at σ=80 (blind dose):  SIFT −0.13  ·  classification +0.08  ·  segmentation +0.19.",
    "Same images, same cleaner — opposite outcomes: denoising hurts texture, helps shape.",
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=18)

# ---------------------------------------------------------------- Slide 13: probes + ranking
s = new_slide("Fine-Tuning — Cost and Generalization")
bullets(s, [
    "Cost of robustness: the fine-tuned classifier drops on clean images, 0.933 → 0.903 (3 points).",
    "Generalization: held-out σ=60 (unseen level) → fine-tuned 0.817 vs. frozen 0.427 — it learned the noise axis, not the training menu.",
    "Robustness ranking (most → least fragile):  SIFT → classification → segmentation.",
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=19)

# ---------------------------------------------------------------- Slide 14: conclusions
s = new_slide("Conclusions & Limitations")
bullets(s, [
    "Fine-tuning is the most reliable recovery, especially under severe degradation.",
    "Blind classical restoration is not free — tuned for the human eye, not the algorithm; dose (over/under) matters.",
    "Carve-out: for shape-based dense tasks, a correctly self-dosed classical cleaner recovers most of the fine-tuning gain, training-free.",
    "Limitations: 300-image val aggregate, binary segmentation, single fine-tune seed, classical restorers only (DL restorers are the next step).",
    ("Repository: " + REPO, 0),
], Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.2), size=18)

OUT.parent.mkdir(exist_ok=True)
prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides)} slides)")
